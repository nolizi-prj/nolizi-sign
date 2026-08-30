"""Generate the small binary fixtures used by test_conversion.py.

Run with the backend venv from the backend/ directory:

    .venv\\Scripts\\python tests\\fixtures\\make_fixtures.py

The generated files are committed to the repo so CI does not need to
regenerate them (and so environments without reportlab/python-docx/openpyxl
installed at arbitrary versions still get byte-identical fixtures).
"""

from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter
from reportlab.pdfgen import canvas

FIXTURES_DIR = Path(__file__).parent


def make_sample_pdf() -> None:
    """Write a 2-page PDF with some text to sample.pdf."""
    path = FIXTURES_DIR / "sample.pdf"
    c = canvas.Canvas(str(path))
    c.drawString(100, 750, "Pumasi Sign sample document - page 1")
    c.showPage()
    c.drawString(100, 750, "Pumasi Sign sample document - page 2")
    c.showPage()
    c.save()


def make_sample_docx() -> None:
    """Write a minimal Word document to sample.docx."""
    path = FIXTURES_DIR / "sample.docx"
    document = Document()
    document.add_paragraph("Pumasi Sign sample document.")
    document.save(str(path))


def make_sample_korean_docx() -> None:
    """Write a Word document whose body is Korean-only text to korean.docx.

    Used to verify converted PDFs embed a CJK-capable font instead of
    rendering tofu boxes (the body deliberately contains no Latin text so
    the embedded-font list reflects only the CJK substitution).
    """
    path = FIXTURES_DIR / "korean.docx"
    document = Document()
    document.add_paragraph("계약서 서명 예시 문서입니다.")
    document.add_paragraph("법인 인감 날인")
    document.save(str(path))


def make_sample_office_fonts_docx() -> None:
    """Write a Word document with explicit Calibri and Cambria runs to
    office-fonts.docx.

    Used to verify converted PDFs keep Office-document metrics: LibreOffice
    must substitute the metric-compatible Carlito/Caladea (or use the real
    fonts where installed), not an arbitrary fallback that reflows the page.
    """
    path = FIXTURES_DIR / "office-fonts.docx"
    document = Document()
    p1 = document.add_paragraph()
    run1 = p1.add_run("Calibri body text, the Office default since 2007.")
    run1.font.name = "Calibri"
    p2 = document.add_paragraph()
    run2 = p2.add_run("Cambria heading text, the Office serif companion.")
    run2.font.name = "Cambria"
    document.save(str(path))


def make_sample_xlsx() -> None:
    """Write a minimal Excel workbook to sample.xlsx."""
    path = FIXTURES_DIR / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = "Pumasi Sign sample document."
    workbook.save(str(path))


def make_sample_wide_xlsx() -> None:
    """Write a 2-sheet workbook (one far wider than a portrait page) to sample-wide.xlsx."""
    path = FIXTURES_DIR / "sample-wide.xlsx"
    workbook = Workbook()
    wide = workbook.active
    wide.title = "Wide"
    for row in range(1, 41):
        for col in range(1, 31):
            wide.cell(row=row, column=col, value=f"Row {row} Col {col}")
    small = workbook.create_sheet("Small")
    small["A1"] = "Second sheet."
    workbook.save(str(path))


def make_sample_pptx() -> None:
    """Write a minimal PowerPoint presentation to sample.pptx."""
    path = FIXTURES_DIR / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    textbox.text_frame.text = "Pumasi Sign sample document."
    presentation.save(str(path))


def make_sample_encrypted_pdf() -> None:
    """Write a password-encrypted PDF to sample-encrypted.pdf, built from sample.pdf."""
    source_path = FIXTURES_DIR / "sample.pdf"
    reader = PdfReader(str(source_path))
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(user_password="secret", owner_password="secret-owner")
    with (FIXTURES_DIR / "sample-encrypted.pdf").open("wb") as f:
        writer.write(f)


if __name__ == "__main__":
    make_sample_pdf()
    make_sample_docx()
    make_sample_korean_docx()
    make_sample_office_fonts_docx()
    make_sample_xlsx()
    make_sample_wide_xlsx()
    make_sample_pptx()
    make_sample_encrypted_pdf()
    print(f"Fixtures written to {FIXTURES_DIR}")
